from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

OUTPUT_FILE = "apresentacao_usina_concreto.pptx"

# Cores
AZUL = RGBColor(22, 54, 92)
LARANJA = RGBColor(230, 126, 34)
CINZA = RGBColor(90, 90, 90)
CINZA_CLARO = RGBColor(230, 233, 238)
BRANCO = RGBColor(255, 255, 255)
PRETO = RGBColor(30, 30, 30)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_bg_white(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BRANCO


def add_top_band(slide, color=AZUL, height=0.55):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    return shape


def add_bottom_band(slide, color=LARANJA, height=0.18):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        prs.slide_height - Inches(height),
        prs.slide_width,
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    return shape


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.8), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.name = "Calibri"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = AZUL

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.72), Inches(1.35), Inches(11.6), Inches(0.4))
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Calibri"
        r2.font.size = Pt(12)
        r2.font.color.rgb = CINZA


def add_bullets(slide, items, left=0.9, top=1.9, width=11.4, height=4.8, font_size=19):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Calibri"
        p.font.size = Pt(font_size)
        p.font.color.rgb = PRETO
        p.space_after = Pt(8)


def add_highlight_box(slide, text, left=0.85, top=5.75, width=11.55, height=0.72):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CINZA_CLARO
    shape.line.color.rgb = AZUL

    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = AZUL


def create_standard_slide(title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(slide)
    add_top_band(slide)
    add_bottom_band(slide)
    add_title(slide, title, subtitle)
    return slide


# Slide 1 - Capa
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg_white(slide)
add_top_band(slide, AZUL, 0.7)
add_bottom_band(slide, LARANJA, 0.22)

title_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(11.3), Inches(1.5))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Fluxo Integrado de Trabalho"
r.font.name = "Calibri"
r.font.size = Pt(28)
r.font.bold = True
r.font.color.rgb = AZUL

p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run()
r2.text = "Usina de Concreto"
r2.font.name = "Calibri"
r2.font.size = Pt(24)
r2.font.bold = True
r2.font.color.rgb = LARANJA

sub_box = slide.shapes.add_textbox(Inches(1.1), Inches(3.25), Inches(11.0), Inches(0.8))
tf2 = sub_box.text_frame
p3 = tf2.paragraphs[0]
p3.alignment = PP_ALIGN.CENTER
r3 = p3.add_run()
r3.text = "Integração entre áreas para eficiência operacional, controle técnico e qualidade do concreto."
r3.font.name = "Calibri"
r3.font.size = Pt(17)
r3.font.color.rgb = CINZA

# Slide 2
slide = create_standard_slide("Objetivo", "Propósito da apresentação")
add_bullets(slide, [
    "Padronizar o fluxo de trabalho entre os setores",
    "Garantir o correto lançamento das informações",
    "Assegurar qualidade na produção e entrega do concreto",
    "Reforçar a importância do controle técnico e operacional",
    "Mostrar o impacto das decisões nos resultados do concreto",
])

# Slide 3
slide = create_standard_slide("Visão Integrada do Processo", "Fluxo macro da operação")
add_bullets(slide, [
    "Comercial: Prospecção → Proposta → Acompanhamento → Pedido no sistema",
    "Técnico: Monitoramento → Conferência → Validação técnica",
    "Operacional: Programação → Produção → Controle de moldagem → Acionamento do laboratório",
    "Gestão da Qualidade: Avaliação de indicadores → Reunião técnica e operacional → Melhoria contínua",
], font_size=17)

# Slide 4
slide = create_standard_slide("Etapa 1: Setor Comercial", "Responsabilidade de origem do processo")
add_bullets(slide, [
    "Prospecção de clientes e oportunidades",
    "Elaboração da proposta comercial",
    "Acompanhamento até a definição do pedido",
    "Lançamento do pedido no sistema Topcon",
    "Registro correto das especificações do concreto",
])
add_highlight_box(slide, "Ponto crítico: o FCK deve ser informado corretamente no lançamento do pedido.")

# Slide 5
slide = create_standard_slide("Qualidade da Informação no Pedido", "Dados críticos para o processo")
add_bullets(slide, [
    "Cliente e obra",
    "Local de entrega",
    "Volume solicitado",
    "Tipo de concreto",
    "FCK especificado",
    "Abatimento / slump",
    "Bombeável ou convencional",
    "Data e horário",
    "Necessidade de moldagem",
    "Observações operacionais relevantes",
], font_size=17)
add_highlight_box(slide, "Informação correta na origem reduz falhas, retrabalho e desvios de qualidade.", top=5.85)

# Slide 6
slide = create_standard_slide("Etapa 2: Departamento Técnico", "Responsabilidade de validação e controle")
add_bullets(slide, [
    "Monitorar as entradas dos pedidos",
    "Conferir a consistência das informações cadastradas",
    "Validar os parâmetros técnicos do concreto solicitado",
    "Identificar desvios, falhas ou inconsistências",
    "Apoiar o alinhamento entre comercial e operação",
    "Preservar os padrões de qualidade da usina",
], font_size=17)
add_highlight_box(slide, "Objetivo central: garantir que o processo produtivo seja iniciado com base em dados confiáveis.", top=5.85)

# Slide 7
slide = create_standard_slide("Etapa 3: Operação / Produção", "Execução com disciplina operacional")
add_bullets(slide, [
    "Receber a demanda programada",
    "Inserir o pedido no fluxo de produção",
    "Organizar carregamento e entrega",
    "Conferir as exigências do pedido",
    "Verificar necessidade de moldagem de corpos de prova",
    "Acionar o laboratório, quando aplicável",
    "Produzir e expedir conforme especificação",
], font_size=17)

# Slide 8
slide = create_standard_slide("Controle Tecnológico e Moldagem", "Rastreabilidade e segurança dos resultados")
add_bullets(slide, [
    "Avaliar exigência contratual do cliente",
    "Confirmar rotina de controle da obra",
    "Aplicar procedimentos internos de qualidade",
    "Acionar o laboratório previamente, quando necessário",
    "Garantir rastreabilidade das amostras",
    "Acompanhar a resistência e o desempenho do concreto",
], font_size=17)

# Slide 9
slide = create_standard_slide("Orientação Operacional Crítica", "Adição de água no concreto")
add_bullets(slide, [
    "A adição de água exige controle rigoroso",
    "Somente pode ocorrer com orientação do encarregado ou gerente da usina",
    "Adição indevida altera o traço e compromete o desempenho",
    "Pode reduzir resistência, afetar abatimento e gerar não conformidade",
], font_size=17)
add_highlight_box(slide, "Regra: nenhuma adição de água sem autorização da liderança da usina.", top=5.85)

# Slide 10
slide = create_standard_slide("Integração entre Técnico e Operacional", "Reunião de desempenho e alinhamento")
add_bullets(slide, [
    "Apresentação dos indicadores de desempenho do concreto",
    "Avaliação do comportamento dos resultados técnicos",
    "Discussão dos impactos das decisões operacionais",
    "Alinhamento entre rotina, produção e qualidade",
    "Definição de ações de melhoria contínua",
], font_size=17)

# Slide 11
slide = create_standard_slide("Indicadores Técnicos de Desempenho", "Números que devem ser monitorados")
add_bullets(slide, [
    "Resistência média",
    "Desvio padrão",
    "Índice de conformidade",
    "Resultados por traço",
    "Resultados por obra",
    "Tendência de desempenho",
    "Ocorrências operacionais relevantes",
    "Registros de desvios e não conformidades",
], font_size=17)

# Slide 12
slide = create_standard_slide("Impacto das Decisões no Resultado Final", "Qualidade é consequência do processo")
add_bullets(slide, [
    "Erro de cadastro do pedido",
    "Falha na especificação do FCK",
    "Comunicação incompleta entre áreas",
    "Produção fora do padrão",
    "Ausência de controle de moldagem",
    "Adição indevida de água",
    "Redução da resistência e aumento do desvio padrão",
    "Retrabalho, perdas e risco reputacional",
], font_size=17)

# Slide 13
slide = create_standard_slide("Compromisso Compartilhado com a Qualidade", "Responsabilidade de toda a equipe")
add_bullets(slide, [
    "Comercial: informação correta na origem",
    "Departamento técnico: validação e monitoramento",
    "Operacional: disciplina na execução",
    "Laboratório: controle e rastreabilidade",
    "Liderança da usina: orientação e tomada de decisão",
    "Motoristas e equipes de campo: cumprimento rigoroso dos procedimentos",
], font_size=17)
add_highlight_box(slide, "Processos bem definidos e responsabilidade compartilhada geram resultados consistentes.", top=5.85)

# Slide 14
slide = create_standard_slide("Conclusão", "Encerramento")
add_bullets(slide, [
    "O fluxo integrado aumenta a eficiência do processo",
    "A qualidade começa no lançamento correto do pedido",
    "O controle técnico reduz riscos e desvios",
    "A operação tem papel decisivo no resultado final",
    "A análise periódica de indicadores fortalece a melhoria contínua",
], font_size=17)
add_highlight_box(slide, "Melhores processos geram melhores concretos, melhores resultados e maior confiança do cliente.", top=5.85)

# Slide 15 - Fechamento
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg_white(slide)
add_top_band(slide, AZUL, 0.7)
add_bottom_band(slide, LARANJA, 0.22)

title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(11.3), Inches(1.0))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Obrigado!"
r.font.name = "Calibri"
r.font.size = Pt(30)
r.font.bold = True
r.font.color.rgb = AZUL

sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(3.1), Inches(11.3), Inches(0.8))
tf2 = sub_box.text_frame
p2 = tf2.paragraphs[0]
p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run()
r2.text = "Compromisso com desempenho, qualidade e evolução contínua."
r2.font.name = "Calibri"
r2.font.size = Pt(18)
r2.font.color.rgb = CINZA

prs.save(OUTPUT_FILE)
print(f"Arquivo gerado com sucesso: {OUTPUT_FILE}")
